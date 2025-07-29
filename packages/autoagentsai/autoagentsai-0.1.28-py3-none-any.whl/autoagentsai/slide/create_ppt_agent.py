from ..client import ChatClient
from ..utils import extract_json
from .util import (
    download_image, 
    parse_markdown_text, 
    apply_inline_formatting, 
    enable_bullet, 
    fill_existing_table, 
    find_nearest_table
)
import os
from typing import Union, List
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

class create_ppt_agent:
    def __init__(self):
        pass

    def outline(self, prompt: str, file_path: Union[str, List[str]]):
        chat_client = ChatClient(
            agent_id="045c418f0dcf4adbb2f15031f06694d1",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
        
        # 处理单文件或多文件情况
        if isinstance(file_path, str):
            files = [file_path]
        else:
            files = file_path
        
        print(f"Debug: 准备处理 {len(files)} 个文件: {files}")
        
        content = ""
        try:
            for event in chat_client.invoke(prompt, files=files):
                if event['type'] == 'start_bubble':
                    print(f"\n{'=' * 20} 消息气泡{event['bubble_id']}开始 {'=' * 20}")
                elif event['type'] == 'token':
                    print(event['content'], end='', flush=True)
                    content += event['content']
                elif event['type'] == 'end_bubble':
                    print(f"\n{'=' * 20} 消息气泡结束 {'=' * 20}")
                elif event['type'] == 'finish':
                    print(f"\n{'=' * 20} 对话完成 {'=' * 20}")
                    break
                elif event['type'] == 'error':
                    print(f"\nDebug: 收到错误事件: {event}")
                    break
                    
        except Exception as e:
            print(f"\nDebug: ChatClient.invoke 发生异常: {type(e).__name__}: {e}")
            # 如果流出现问题，返回错误信息而不是空字符串
            if not content.strip():
                content = f"Stream error: {str(e)}"
        
        print(f"\nDebug: 最终返回内容长度: {len(content)}")
        content = extract_json(content)
        return content

    # def cover(self):
    #     pass

    # def content(self):
    #     pass

    # def conclusion(self):
    #     pass
    
    # def save(self, file_path: str):
    #     pass


    def fill(self, data: dict, template_file_path: str = "template-1.pptx", output_file_path: str = "output-1.pptx"):
        # 加载 PPTX 模板
        prs = Presentation(template_file_path)

        # 用于存储需要清理的临时文件
        temp_files = []

        # 处理远程图片下载
        processed_data = {}
        # 支持的图片文件后缀
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg')
        
        for key, value in data.items():
            # 检查是否是远程图片URL
            if isinstance(value, str) and value.startswith(('http://', 'https://')):
                # 检查URL是否以图片文件后缀结尾（忽略查询参数）
                url_path = value.split('?')[0].lower()  # 去掉查询参数并转为小写
                if url_path.endswith(image_extensions):
                    # 尝试下载远程图片
                    local_image_path = download_image(value)
                    if local_image_path:
                        temp_files.append(local_image_path)
                        processed_data[key] = local_image_path
                        print(f"成功下载图片: {key} = {value} -> {local_image_path}")
                    else:
                        print(f"跳过下载失败的图片: {key} = {value}")
                        # 下载失败时不添加到processed_data，这样图片占位符就会被跳过
                else:
                    print(f"跳过非图片URL: {key} = {value} (不支持的文件类型)")
                    processed_data[key] = value  # 保留原值，可能是本地路径
            else:
                processed_data[key] = value

        # 1. 表格填充
        table_requests = []  # [(占位符形状, key, data)]
        all_tables = []  # 所有表格形状
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # 收集表格占位符
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text.startswith("{{#") and text.endswith("}}"):
                        key = text[3:-2].strip()  # 去掉 {{# 和 }}
                        if key in processed_data:
                            table_requests.append((shape, key, processed_data[key]))
                
                # 收集所有表格
                if shape.has_table:
                    all_tables.append(shape)
        
        # 为每个表格占位符找到最近的表格并填充
        shapes_to_remove = []
        processed_tables = set()
        
        for placeholder_shape, key, table_data in table_requests:
            available_tables = [t for t in all_tables if id(t) not in processed_tables]
            if not available_tables:
                available_tables = all_tables
            
            nearest_table_shape = find_nearest_table(placeholder_shape, available_tables)
            if nearest_table_shape:
                print(f"占位符 '{{#{key}}}' 匹配到最近的表格")
                fill_existing_table(nearest_table_shape.table, table_data)
                processed_tables.add(id(nearest_table_shape))
            
            shapes_to_remove.append(placeholder_shape)
        
        # 删除表格占位符文本框
        for shape in shapes_to_remove:
            shape._element.getparent().remove(shape._element)
        
        # 2. 文本、图片填充
        for slide in prs.slides:
            for shape in list(slide.shapes):  # list() to allow removal
                if not shape.has_text_frame:
                    continue
            
                text = shape.text.strip()
                if text.startswith("{{") and text.endswith("}}"):
                    key = text[2:-2].strip()  # 去掉 {{}}
                    content_type = "text"

                    # 判断类型前缀
                    if key.startswith("@"):
                        key = key[1:]
                        content_type = "image"
                    elif key.startswith("#"):
                        # 表格已经在上面处理过了，跳过
                        continue

                    value = processed_data.get(key)
                    if value is None:
                        continue

                    if content_type == "text":
                        # 检查是否包含Markdown格式
                        if isinstance(value, str) and any(marker in value for marker in ['*', '#', '`', '\n']):
                            # 使用Markdown解析
                            parse_markdown_text(shape.text_frame, value)
                        elif isinstance(value, list):
                            # 处理列表数据，每项作为bullet point
                            tf = shape.text_frame
                            tf.clear()
                            for i, item in enumerate(value):
                                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                                if isinstance(item, str) and any(marker in item for marker in ['*', '#', '`']):
                                    apply_inline_formatting(p, item)
                                else:
                                    p.text = str(item)
                                    p.font.size = Pt(14)
                                    p.alignment = PP_ALIGN.LEFT
                                    enable_bullet(p)
                        else:
                            # 普通文本
                            shape.text_frame.text = str(value)

                    elif content_type == "image":
                        # 获取位置并删除原文本框
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes._spTree.remove(shape._element)
                            
                        # 确保图片路径存在
                        if os.path.exists(value):
                            slide.shapes.add_picture(value, left, top, width=width, height=height)
                            print(f"成功替换图片: {key}")
                    else:
                            print(f"警告: 图片文件不存在: {value}")

        # 保存为新PPT
        prs.save(output_file_path)
        
        # 清理临时文件
        for temp_file in temp_files:
            try:
                os.unlink(temp_file)
                print(f"清理临时文件: {temp_file}")
            except Exception as e:
                print(f"清理临时文件失败: {temp_file}, 错误: {e}")