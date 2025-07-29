from pptx import Presentation

def handle_pptx(path, output_format="text"):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text.strip())
        slides.append("\n".join(content))

    text = "\n---\n".join(slides)
    return text if output_format == "text" else {"slides": slides}
